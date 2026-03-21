"""Извлечение текста из загруженных файлов (PDF, DOCX, PPTX, TXT)."""

from pathlib import Path

from app.services.text_sanitize import sanitize_postgres_text


def extract_text_from_file(mime_type: str, path: Path) -> str:
    mt = mime_type.split(";")[0].strip().lower()
    if mt == "application/pdf":
        raw = _pdf(path)
    elif mt == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        raw = _docx(path)
    elif mt == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        raw = _pptx(path)
    elif mt == "text/plain":
        raw = _txt(path)
    else:
        raise ValueError(f"Неподдерживаемый тип: {mime_type}")
    return sanitize_postgres_text(raw)


def _pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            parts.append(t)
    return "\n".join(parts)


def _docx(path: Path) -> str:
    from docx import Document

    document = Document(str(path))
    return "\n".join(p.text for p in document.paragraphs if p.text)


def _pptx(path: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"\n--- Слайд {i} ---\n")
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shape.table
                for row in tbl.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
            elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
    return "\n".join(p for p in parts if p)


def _txt(path: Path) -> str:
    raw = path.read_bytes()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")
